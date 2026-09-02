#!/usr/bin/env python3
"""
Build the destination.com media-kit slides as shareable JPGs.

    python3 media-kit/build.py

Reads media-kit/config.json, renders one HTML page per slide, screenshots each
with headless Chromium, and writes JPGs to media-kit/out/.

Any config value left as null renders as a visible "ADD DATA" chip, and the
slides that depend on data carry a TEMPLATE banner until every figure on them
is populated. Nothing is invented at render time: every number shown comes
from config.json.

config.json's `figures_basis` says how the numbers it holds were arrived at.
While it reads "estimated", every data slide carries an amber ESTIMATES pill
and prints `basis_note` in its footer, so a modelled figure is never presented
as a measured one. Set it to "measured" only once the figures really are.
"""

import json
import os
import shutil
import subprocess
import sys

from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__))
WORK = os.path.join(ROOT, "_work")
OUT = os.path.join(ROOT, "out")
FONT_CSS = os.path.join(ROOT, "fonts", "fonts-embedded.css")

WIDTH, HEIGHT = 1600, 900
SCALE = 2
FINAL_WIDTH = 2400
JPG_QUALITY = 90

# headless_shell is preferred: --window-size maps 1:1 onto the viewport. The full Chrome
# binary reserves ~87px of window for browser chrome even under --headless=new, which
# silently shrinks the viewport and lets slide content fall off the bottom of the capture.
CHROME_CANDIDATES = [
    "/opt/pw-browsers/chromium_headless_shell-1194/chrome-linux/headless_shell",
    "/opt/pw-browsers/chromium-1194/chrome-linux/chrome",
    shutil.which("chromium"),
    shutil.which("chromium-browser"),
    shutil.which("google-chrome"),
]


def find_chrome():
    for c in CHROME_CANDIDATES:
        if c and os.path.exists(c):
            return c
    import glob

    hits = (glob.glob("/opt/pw-browsers/chromium_headless_shell*/chrome-linux/headless_shell")
            + glob.glob("/opt/pw-browsers/chromium*/chrome-linux/chrome"))
    if hits:
        return sorted(hits)[-1]
    sys.exit("No Chromium binary found. Set CHROME env var or install chromium.")


CHROME = os.environ.get("CHROME") or find_chrome()
IS_SHELL = os.path.basename(CHROME) == "headless_shell"

# ---------------------------------------------------------------- brand tokens

CSS = """
:root {
  --primary:#1B6CA8; --primary-dark:#15557F; --secondary:#0D9B76;
  --accent:#C4703E; --gold:#F4A124;
  --cream:#FAF7F2; --cream-dark:#F0EBE3; --white:#FFFFFF;
  --ink:#1A1A1A; --ink-mid:#4A4A4A; --ink-light:#7A7A7A;
  --border:#E5E0D8;
}
* { margin:0; padding:0; box-sizing:border-box; }
html,body { width:1600px; height:900px; overflow:hidden; }
body {
  font-family:'DM Sans', 'Liberation Sans', sans-serif;
  color:var(--ink); background:var(--white);
  -webkit-font-smoothing:antialiased;
}
h1,h2,h3,h4 { font-family:'Playfair Display','Liberation Serif',serif; font-weight:600; }

.slide { width:1600px; height:900px; position:relative; display:flex; flex-direction:column; }

/* ---- slide chrome ---- */
.shead { padding:38px 56px 22px; display:flex; align-items:flex-end; justify-content:space-between;
         border-bottom:1px solid var(--border); }
.shead .num { font-size:11px; letter-spacing:2.2px; text-transform:uppercase;
              color:var(--accent); font-weight:700; margin-bottom:7px; }
.shead h2 { font-size:33px; line-height:1.1; }
.shead .sub { font-size:14px; color:var(--ink-mid); margin-top:9px; max-width:720px; line-height:1.5; }
.shead .brandmark { font-family:'Playfair Display',serif; font-size:15px; color:var(--ink-light);
                    letter-spacing:.5px; white-space:nowrap; padding-bottom:4px; }
.sbody { flex:1; display:flex; gap:34px; padding:30px 56px 26px; min-height:0; }
.sfoot { padding:0 56px 22px; font-size:11px; color:var(--ink-light); line-height:1.5; }

/* ---- placeholder ---- */
.ph { display:inline-block; font-family:'DM Sans',sans-serif; font-size:11px; font-weight:700;
      letter-spacing:1.3px; color:#B9481F; background:repeating-linear-gradient(45deg,
      #FDF0E8, #FDF0E8 6px, #F9E3D5 6px, #F9E3D5 12px);
      border:1.5px dashed #E0885F; border-radius:5px; padding:4px 10px; vertical-align:middle; }
.ph-lg { font-size:15px; padding:8px 14px; letter-spacing:1.6px; }
.tmpl-banner { background:#B9481F; color:#fff; font-size:11.5px; font-weight:700; letter-spacing:1.6px;
      text-transform:uppercase; padding:9px 56px; display:flex; gap:12px; align-items:center; }
.tmpl-banner span.dot { width:7px; height:7px; border-radius:50%; background:#F4A124; display:inline-block; }
.est-pill { display:inline-flex; align-items:center; gap:7px; background:#FDF6E8; border:1px solid #E3C489;
      color:#7A5000; font-size:9.5px; font-weight:700; letter-spacing:1.4px; text-transform:uppercase;
      padding:6px 13px; border-radius:20px; white-space:nowrap; }
.est-pill span.d { width:6px; height:6px; border-radius:50%; background:#D9922B; display:inline-block; }
.shead .rt { display:flex; flex-direction:column; align-items:flex-end; gap:10px; }

/* ---- stat tiles ---- */
.tiles { display:flex; gap:16px; }
.tile { flex:1; background:var(--cream); border:1px solid var(--border); border-radius:12px;
        padding:22px 22px 20px; display:flex; flex-direction:column; justify-content:space-between; }
.tile .k { font-size:10.5px; letter-spacing:1.5px; text-transform:uppercase; color:var(--ink-light);
           font-weight:600; line-height:1.45; }
.tile .v { font-family:'Playfair Display',serif; font-size:42px; font-weight:600; color:var(--primary);
           line-height:1; margin-top:14px; }
.tile .v small { font-size:17px; font-weight:600; color:var(--ink-mid); margin-left:3px;
                 font-family:'DM Sans',sans-serif; }

/* ---- bar charts ---- */
.chart { display:flex; flex-direction:column; gap:11px; }
.chart .ct { font-size:11px; letter-spacing:1.6px; text-transform:uppercase; color:var(--ink-light);
             font-weight:700; margin-bottom:3px; }
.row { display:grid; grid-template-columns:118px 1fr 56px; align-items:center; gap:14px; }
.row .lab { font-size:13px; color:var(--ink-mid); text-align:right; font-weight:500; }
.row .track { height:20px; background:var(--cream-dark); border-radius:4px; position:relative; overflow:hidden; }
.row .track.empty { background:repeating-linear-gradient(45deg,#F6F2EC,#F6F2EC 5px,#EEE7DD 5px,#EEE7DD 10px);
                    border:1px dashed #D8CDBC; }
.row .fill { height:100%; background:var(--primary); border-radius:0 4px 4px 0; }
.row .val { font-size:13.5px; font-weight:700; color:var(--ink); }
.row .val.na { font-size:9.5px; font-weight:700; color:#B9481F; letter-spacing:.8px; }
.chart .rows { display:flex; flex-direction:column; gap:11px; }
.chart.fill { flex:1; min-height:0; }
.chart.fill .rows { flex:1; justify-content:center; gap:20px; }
.chart.fill .track { height:27px; }
.stretch { height:100%; display:flex; flex-direction:column; justify-content:space-evenly; }

/* ---- split bar ---- */
.split { display:flex; height:44px; border-radius:8px; overflow:hidden; gap:2px; background:var(--white); }
.split .seg { display:flex; align-items:center; justify-content:center; color:#fff; font-size:13px;
              font-weight:700; letter-spacing:.3px; }
.split .seg.empty { background:repeating-linear-gradient(45deg,#F6F2EC,#F6F2EC 5px,#EEE7DD 5px,#EEE7DD 10px);
                    border:1px dashed #D8CDBC; color:#B9481F; font-size:10px; letter-spacing:1.2px; }
.legend { display:flex; gap:18px; margin-top:10px; flex-wrap:wrap; }
.legend .li { display:flex; align-items:center; gap:7px; font-size:12px; color:var(--ink-mid); }
.legend .sw { width:11px; height:11px; border-radius:3px; }

/* ---- device frames ---- */
.frame { background:var(--white); border:1px solid var(--border); border-radius:11px;
         box-shadow:0 10px 34px rgba(0,0,0,.10); overflow:hidden; display:flex; flex-direction:column; }
.chrome { height:34px; background:var(--cream-dark); display:flex; align-items:center; gap:7px;
          padding:0 13px; border-bottom:1px solid var(--border); flex-shrink:0; }
.chrome i { width:9px; height:9px; border-radius:50%; background:#CFC6B8; display:block; }
.chrome .url { flex:1; height:19px; background:var(--white); border-radius:10px; margin-left:9px;
               font-size:10px; color:var(--ink-light); display:flex; align-items:center;
               padding:0 11px; font-family:'DM Sans',sans-serif; }

/* ---- article mockup ---- */
.pg { flex:1; overflow:hidden; background:var(--white); }
.nav { display:flex; align-items:center; justify-content:space-between; padding:11px 20px;
       border-bottom:1px solid var(--border); }
.nav .logo { font-family:'Playfair Display',serif; font-size:14px; }
.nav .lk { display:flex; gap:14px; }
.nav .lk span { font-size:7.5px; letter-spacing:1.1px; text-transform:uppercase; color:var(--ink-mid); }
.nav .cta { background:var(--accent); color:#fff; font-size:7.5px; letter-spacing:.9px;
            text-transform:uppercase; padding:6px 11px; border-radius:3px; font-weight:700; }
.crumb { padding:11px 20px 0; font-size:8.5px; color:var(--ink-light); letter-spacing:.4px; }
.art { padding:8px 20px 0; }
.art h1 { font-size:23px; line-height:1.18; letter-spacing:-.2px; }
.art .byl { font-size:9px; color:var(--ink-light); margin-top:9px; padding-bottom:12px;
            border-bottom:1px solid var(--border); }
.hero { height:118px; margin:13px 20px 0; border-radius:7px;
        background:linear-gradient(135deg,#1B6CA8 0%,#2E8BB8 42%,#0D9B76 100%); position:relative; }
.hero:after { content:'SS163 · Amalfi Coast'; position:absolute; left:13px; bottom:10px; color:#fff;
              font-size:8.5px; letter-spacing:1.3px; text-transform:uppercase; opacity:.9; font-weight:600; }
.p { font-size:9.5px; line-height:1.75; color:var(--ink-mid); padding:0 20px; margin-top:11px; }
.p b { color:var(--ink); }
.h2 { font-family:'Playfair Display',serif; font-size:14px; padding:0 20px; margin-top:16px; }

/* ---- ad slot ---- */
.slot { margin:15px 20px; border:2.5px solid var(--accent); border-radius:7px; background:#FCF4EE;
        display:flex; flex-direction:column; align-items:center; justify-content:center; position:relative; }
.slot .tag { position:absolute; top:-10px; left:13px; background:var(--accent); color:#fff;
             font-size:8px; font-weight:700; letter-spacing:1.3px; text-transform:uppercase;
             padding:3px 9px; border-radius:3px; }
.slot .sz { font-size:13px; font-weight:700; color:var(--accent); letter-spacing:.5px; }
.slot .sz small { display:block; font-size:8.5px; font-weight:600; color:var(--ink-light);
                  letter-spacing:1.1px; text-transform:uppercase; margin-top:4px; }

/* ---- spec panel ---- */
.spec { width:392px; flex-shrink:0; display:flex; flex-direction:column; gap:15px; }
.spec .card { background:var(--cream); border:1px solid var(--border); border-radius:12px; padding:19px 21px; }
.spec .card.hi { background:#FCF4EE; border-color:#E8C4AA; }
.spec .ck { font-size:10px; letter-spacing:1.7px; text-transform:uppercase; color:var(--accent);
            font-weight:700; margin-bottom:9px; }
.spec .cv { font-size:13.5px; line-height:1.6; color:var(--ink); }
.spec .cv.mid { color:var(--ink-mid); font-size:12.5px; }
.spec dl { display:flex; flex-direction:column; gap:12px; }
.spec dt { font-size:9.5px; letter-spacing:1.5px; text-transform:uppercase; color:var(--ink-light);
           font-weight:700; margin-bottom:3px; }
.spec dd { font-size:13px; line-height:1.55; color:var(--ink); }

/* ---- native cards ---- */
.ncards { display:flex; gap:9px; padding:0 20px; margin-top:11px; }
.ncard { flex:1; border:1px solid var(--border); border-radius:6px; overflow:hidden; background:#fff; }
.ncard .im { height:52px; background:linear-gradient(135deg,#D8E4EC,#C3D6E2); }
.ncard .bd { padding:8px 9px 10px; }
.ncard .nm { font-size:9px; font-weight:700; color:var(--ink); }
.ncard .mt { font-size:7.5px; color:var(--ink-light); margin-top:3px; }
.ncard .pr { font-size:9px; font-weight:700; color:var(--secondary); margin-top:6px; }
.ncard.sp { border:2.5px solid var(--accent); position:relative; }
.ncard.sp .im { background:linear-gradient(135deg,#1B6CA8,#0D9B76); }
.ncard.sp .flag { position:absolute; top:5px; right:5px; background:var(--accent); color:#fff;
                  font-size:6.5px; font-weight:700; letter-spacing:.9px; text-transform:uppercase;
                  padding:2px 6px; border-radius:2px; }

/* ---- phone ---- */
.phone { width:266px; border:9px solid #23272B; border-radius:30px; overflow:hidden;
         box-shadow:0 12px 36px rgba(0,0,0,.16); background:#fff; display:flex; flex-direction:column;
         flex-shrink:0; }
.phone .notch { height:19px; background:#23272B; display:flex; align-items:center; justify-content:center; }
.phone .notch i { width:56px; height:5px; border-radius:3px; background:#3A4046; display:block; }

/* ---- email ---- */
.mail { flex:1; background:var(--cream); padding:15px; overflow:hidden; }
.mail .card { background:#fff; border:1px solid var(--border); border-radius:8px; overflow:hidden; }
.mail .mh { padding:13px 17px; border-bottom:1px solid var(--border); }
.mail .mh .t { font-family:'Playfair Display',serif; font-size:15px; }
.mail .mh .d { font-size:8px; color:var(--ink-light); letter-spacing:1.2px; text-transform:uppercase;
               margin-top:5px; }
.mail .item { padding:11px 17px; border-bottom:1px solid var(--border); }
.mail .item .it { font-size:10.5px; font-weight:700; }
.mail .item .ix { font-size:8.5px; color:var(--ink-light); margin-top:3px; line-height:1.5; }

/* ---- table ---- */
table.spectbl { width:100%; border-collapse:collapse; }
table.spectbl th { font-size:10px; letter-spacing:1.6px; text-transform:uppercase; color:var(--ink-light);
                   font-weight:700; text-align:left; padding:0 16px 11px; border-bottom:1.5px solid var(--border); }
table.spectbl td { font-size:13px; color:var(--ink-mid); padding:15px 16px; border-bottom:1px solid var(--border);
                   vertical-align:top; line-height:1.5; }
table.spectbl td.nm { color:var(--ink); font-weight:700; font-size:13.5px; width:212px; }
table.spectbl td.sz { font-family:'Liberation Mono',monospace; font-size:11.5px; color:var(--primary);
                      width:238px; }
"""


# ---------------------------------------------------------------- helpers

def load_cfg():
    with open(os.path.join(ROOT, "config.json")) as f:
        return json.load(f)


def fmt_int(v):
    return f"{v:,}" if isinstance(v, (int, float)) else None


def chip(size=""):
    cls = "ph ph-lg" if size == "lg" else "ph"
    return f'<span class="{cls}">ADD DATA</span>'


def tile(label, value, unit=""):
    if value is None:
        v = chip("lg")
    else:
        shown = fmt_int(value) if unit != "%" else f"{value}"
        v = f'{shown}<small>{unit}</small>' if unit else shown
    return f'<div class="tile"><div class="k">{label}</div><div class="v">{v}</div></div>'


def bars(title, items, maxpct=None, fill=False, label_width=None):
    vals = [i["pct"] for i in items if i.get("pct") is not None]
    top = maxpct or (max(vals) if vals else 100)
    top = max(top, 1)
    # A long category name wraps to two lines in the default gutter and knocks the
    # row off the baseline its neighbours sit on; widen the gutter instead.
    style = f' style="grid-template-columns:{label_width}px 1fr 56px"' if label_width else ""
    rows = []
    for it in items:
        p = it.get("pct")
        if p is None:
            rows.append(
                f'<div class="row"{style}><div class="lab">{it["label"]}</div>'
                f'<div class="track empty"></div><div class="val na">ADD</div></div>'
            )
        else:
            w = max(2.0, p / top * 100)
            rows.append(
                f'<div class="row"{style}><div class="lab">{it["label"]}</div>'
                f'<div class="track"><div class="fill" style="width:{w:.1f}%"></div></div>'
                f'<div class="val">{p}%</div></div>'
            )
    cls = "chart fill" if fill else "chart"
    return (f'<div class="{cls}"><div class="ct">{title}</div>'
            f'<div class="rows">{"".join(rows)}</div></div>')


SPLIT_COLORS = ["#1B6CA8", "#C4703E", "#0D9B76"]


def split_bar(title, items):
    vals = [i.get("pct") for i in items]
    if any(v is None for v in vals):
        segs = '<div class="seg empty" style="flex:1">ADD DATA</div>'
        leg = ""
    else:
        # A segment under ~9% cannot hold "Label NN%" without clipping. Drop the inline
        # label there and let the legend carry identity (never color-alone: the legend
        # is always rendered for a split bar).
        segs = "".join(
            f'<div class="seg" style="flex:{v};background:{SPLIT_COLORS[i % 3]}">'
            f'{(items[i]["label"] + " " if v >= 9 else "")}{v}%</div>'
            for i, v in enumerate(vals)
        )
        leg = '<div class="legend">' + "".join(
            f'<div class="li"><span class="sw" style="background:{SPLIT_COLORS[i % 3]}"></span>'
            f'{items[i]["label"]}</div>' for i in range(len(items))
        ) + "</div>"
    return (f'<div class="chart"><div class="ct">{title}</div>'
            f'<div class="split">{segs}</div>{leg}</div>')


def estimate_pill(cfg):
    """Amber 'estimates' pill for slides whose figures are not yet measured."""
    if cfg.get("figures_basis") != "estimated":
        return ""
    return ('<div class="est-pill"><span class="d"></span>'
            'Estimates &middot; pending verification</div>')


def head(num, title, sub, brand, pill=""):
    return (f'<div class="shead"><div><div class="num">{num}</div><h2>{title}</h2>'
            f'<div class="sub">{sub}</div></div>'
            f'<div class="rt">{pill}<div class="brandmark">{brand["property"]}</div></div></div>')


def banner_if_missing(*groups):
    """Render the TEMPLATE banner if any figure on the slide is still null."""
    missing = False
    for g in groups:
        if g is None:
            missing = True
        elif isinstance(g, list):
            for it in g:
                if isinstance(it, dict) and (it.get("pct") is None and it.get("value") is None):
                    missing = True
    if not missing:
        return ""
    return ('<div class="tmpl-banner"><span class="dot"></span>'
            'Template — figures not yet populated. Do not send until every field is filled '
            'from GA4 and the ESP.</div>')


# ---------------------------------------------------------------- slides

def slide_cover(c):
    b, r = c["brand"], c["reach"]
    return f"""
<div class="slide" style="background:linear-gradient(115deg,#FAF7F2 0%,#F3EDE4 52%,#EFE8DC 100%);">
  <div style="position:absolute;right:-120px;top:-120px;width:520px;height:520px;border-radius:50%;
       background:linear-gradient(135deg,rgba(27,108,168,.10),rgba(13,155,118,.10));"></div>
  <div style="padding:78px 56px 0;position:relative;">
    <div style="font-size:11px;letter-spacing:3px;text-transform:uppercase;color:var(--accent);
         font-weight:700;">{b['kicker']}</div>
    <h1 style="font-size:70px;line-height:1.02;margin-top:20px;letter-spacing:-1.2px;">
      {b['property']}</h1>
    <div style="font-size:20px;color:var(--ink-mid);margin-top:18px;max-width:660px;line-height:1.5;">
      {b['tagline']}</div>
    <div style="width:62px;height:3px;background:var(--accent);margin:30px 0 0;"></div>
    <div style="font-size:13px;color:var(--ink-light);margin-top:22px;letter-spacing:.4px;">
      Prepared for <b style="color:var(--ink);">{b['prepared_for']}</b> · {b['date']}</div>
  </div>
  <div style="margin-top:auto;padding:0 56px 34px;position:relative;">
    <div class="tiles">
      {tile("Monthly readers", r["monthly_readers"])}
      {tile("Newsletter subscribers", r["newsletter_subscribers"])}
      {tile("Avg. engaged time", r["avg_engaged_time_min"], "min")}
      {tile("Returning readers", r["returning_reader_pct"], "%")}
    </div>
    <div style="display:flex;justify-content:space-between;align-items:center;margin-top:26px;
         font-size:12px;color:var(--ink-light);">
      <span>{b['contact_name']} · {b['contact_email']}</span>
      {estimate_pill(c) or '<span>Placement specifications and audience detail follow</span>'}
    </div>
  </div>
</div>"""


ARTICLE_TOP = """
<div class="nav"><div class="logo">destination</div>
  <div class="lk"><span>Europe</span><span>Asia</span><span>Americas</span><span>Guides</span>
    <span>Points &amp; Miles</span></div>
  <div class="cta">Plan a trip</div></div>
<div class="crumb">Home &rsaquo; Europe &rsaquo; Italy &rsaquo; Amalfi Coast</div>
<div class="art">
  <h1>The Amalfi Coast Road Trip: Route, Stops, Parking Reality &amp; What to Skip</h1>
  <div class="byl">By Elena Moretti · 28 August 2026 · 9 min read</div>
</div>
<div class="hero"></div>
<div class="p"><b>Rental cars are banned from the coast road on alternating days through the
  summer.</b> The <i>targhe alterne</i> system has run since 2022, it applies specifically to
  hire cars, and people still arrive not knowing about it.</div>
<div class="h2">The Quick Answer</div>
<div class="p">On odd-numbered dates, odd-numbered plates are barred; on even dates, even plates
  are barred. Roughly June–September, 10:00–18:00, along the full SS163 from Vietri sul Mare to
  Positano. Parking is the real constraint, not the driving.</div>
"""


def slide_in_article(c):
    p = c["placements"][0]
    return f"""
<div class="slide">
  {head("Placement 01", "In-article display", "The standard display unit, placed after the second section of every guide — where the reader has committed to the article and is still forming the trip.", c["brand"])}
  <div class="sbody">
    <div class="frame" style="flex:1;">
      <div class="chrome"><i></i><i></i><i></i>
        <div class="url">destination.com/europe/italy/amalfi-coast-road-trip</div></div>
      <div class="pg">
        {ARTICLE_TOP}
        <div class="slot" style="height:118px;"><span class="tag">Ad placement</span>
          <span class="sz">970 × 250<small>Billboard · desktop</small></span></div>
        <div class="p">For most people the ferry is the better answer between May and September.</div>
      </div>
    </div>
    <div class="spec">
      <div class="card hi"><div class="ck">Unit</div><div class="cv">{p['name']}</div></div>
      <div class="card"><dl>
        <div><dt>Sizes accepted</dt><dd>{p['sizes']}</dd></div>
        <div><dt>Position</dt><dd>{p['position']}</dd></div>
        <div><dt>Delivery</dt><dd>{p['note']}</dd></div>
      </dl></div>
      <div class="card"><div class="ck">Why it works here</div>
        <div class="cv mid">The reader arrived from a search for a specific destination.
        By this scroll depth they have read the constraint that shapes the trip and are
        deciding how to book around it.</div></div>
    </div>
  </div>
  <div class="sfoot">Mockup for placement illustration. Creative shown at relative scale;
    served sizes are as specified above.</div>
</div>"""


def slide_native(c):
    p = c["placements"][1]
    return f"""
<div class="slide">
  {head("Placement 02", "&ldquo;Where to Stay&rdquo; native module", "A sponsored card inside the accommodation section of each destination guide. The highest-intent unit on the site — the destination is already chosen.", c["brand"])}
  <div class="sbody">
    <div class="frame" style="flex:1;">
      <div class="chrome"><i></i><i></i><i></i>
        <div class="url">destination.com/europe/italy/amalfi-coast-road-trip#where-to-stay</div></div>
      <div class="pg">
        <div class="nav"><div class="logo">destination</div>
          <div class="lk"><span>Europe</span><span>Asia</span><span>Americas</span><span>Guides</span>
            <span>Points &amp; Miles</span></div>
          <div class="cta">Plan a trip</div></div>
        <div class="h2" style="margin-top:20px;">Where to Stay on the Amalfi Coast</div>
        <div class="p">Positano is the postcard and prices it accordingly. Praiano sits between
          Positano and Amalfi with the same coastline and materially lower rates. Salerno is the
          budget base with the best transport links.</div>
        <div class="ncards">
          <div class="ncard"><div class="im"></div><div class="bd">
            <div class="nm">Hotel Marincanto</div><div class="mt">Positano · 4★ · sea view</div>
            <div class="pr">from €340 / night</div></div></div>
          <div class="ncard sp"><span class="flag">Sponsored</span><div class="im"></div>
            <div class="bd"><div class="nm">Partner placement</div>
            <div class="mt">Logo · headline · 12-word description</div>
            <div class="pr">Advertiser CTA &rsaquo;</div></div></div>
          <div class="ncard"><div class="im"></div><div class="bd">
            <div class="nm">Casa Angelina</div><div class="mt">Praiano · 5★ · clifftop</div>
            <div class="pr">from €410 / night</div></div></div>
        </div>
        <div class="slot" style="height:40px;margin:14px 20px;"><span class="tag">Native slot</span>
          <span class="sz" style="font-size:10.5px;">Card 2 of 3 · responsive · matches editorial styling</span></div>
        <div class="h2" style="margin-top:6px;">Getting Between Towns</div>
        <div class="p">The ferry runs Salerno–Amalfi–Positano from April to October and is faster
          than the road in peak season. SITA buses cover the same route year-round for a few euro,
          but they fill at the Amalfi terminus and you may not get on the first one.</div>
        <div class="p">If you are staying in Praiano, check whether your hotel runs a shuttle to
          the SS163 — the walk up from the shore is longer than it looks on a map.</div>
      </div>
    </div>
    <div class="spec">
      <div class="card hi"><div class="ck">Unit</div><div class="cv">{p['name']}</div></div>
      <div class="card"><dl>
        <div><dt>Format</dt><dd>{p['sizes']}</dd></div>
        <div><dt>Position</dt><dd>{p['position']}</dd></div>
        <div><dt>Labelling</dt><dd>Marked &ldquo;Sponsored&rdquo; on the card, per FTC disclosure
          requirements. Editorial cards either side are unpaid.</dd></div>
      </dl></div>
      <div class="card"><div class="ck">Why it works here</div>
        <div class="cv mid">{p['note']}. The unit sits in the decision itself rather than
        beside it.</div></div>
    </div>
  </div>
  <div class="sfoot">Editorial properties shown are illustrative. Sponsored cards are always
    labelled and never displace an editorial recommendation.</div>
</div>"""


def slide_mobile(c):
    p = c["placements"][0]
    return f"""
<div class="slide">
  {head("Placement 03", "Mobile in-feed", "The majority of reading happens on a phone. The mobile unit is a single 300&thinsp;&times;&thinsp;250 in the article flow — one unit per viewport, never a stack.", c["brand"])}
  <div class="sbody" style="align-items:flex-start;">
    <div style="flex:1;display:flex;justify-content:center;gap:34px;align-items:flex-start;">
      <div class="phone" style="height:588px;">
        <div class="notch"><i></i></div>
        <div class="pg">
          <div class="nav" style="padding:9px 13px;"><div class="logo" style="font-size:12px;">destination</div>
            <div class="cta" style="font-size:6.5px;padding:5px 9px;">Plan a trip</div></div>
          <div class="art" style="padding:10px 14px 0;">
            <h1 style="font-size:16px;">The Amalfi Coast Road Trip: Route, Parking &amp; What to Skip</h1>
            <div class="byl" style="font-size:8px;">By Elena Moretti · 9 min read</div></div>
          <div class="hero" style="height:82px;margin:11px 14px 0;"></div>
          <div class="p" style="padding:0 14px;font-size:8.5px;"><b>Rental cars are banned from the
            coast road on alternating days through the summer.</b> The <i>targhe alterne</i> system
            has run since 2022 and applies specifically to hire cars.</div>
          <div class="slot" style="height:172px;margin:13px 14px;">
            <span class="tag">Ad placement</span>
            <span class="sz">300 × 250<small>Mobile MPU</small></span></div>
          <div class="p" style="padding:0 14px;font-size:8.5px;">Parking is the real constraint,
            not the driving.</div>
        </div>
      </div>
      <div style="width:250px;padding-top:26px;">
        <div style="font-size:11px;letter-spacing:1.7px;text-transform:uppercase;color:var(--accent);
             font-weight:700;">Reading behaviour</div>
        <div style="font-size:13px;color:var(--ink-mid);line-height:1.65;margin-top:12px;">
          Mobile share of sessions and engaged time are on the audience slides. The unit is
          lazy-loaded below the fold and renders only when it enters the viewport, so served
          impressions and viewable impressions stay close together.</div>
        <div style="margin-top:22px;padding-top:18px;border-top:1px solid var(--border);
             font-size:13px;color:var(--ink-mid);line-height:1.65;">
          No interstitials, no anchor units, no auto-refresh. One in-view unit at a time.</div>
      </div>
    </div>
    <div class="spec">
      <div class="card hi"><div class="ck">Unit</div><div class="cv">In-article display — mobile</div></div>
      <div class="card"><dl>
        <div><dt>Size</dt><dd>300 × 250</dd></div>
        <div><dt>Position</dt><dd>{p['position']}</dd></div>
        <div><dt>Density</dt><dd>{p['note']}</dd></div>
      </dl></div>
      <div class="card"><div class="ck">Sidebar (desktop only)</div>
        <div class="cv mid">A 300&thinsp;&times;&thinsp;600 sticky rail unit runs alongside the
        article body on desktop, holding share of voice for the length of the read.</div></div>
    </div>
  </div>
  <div class="sfoot">Mockup for placement illustration. Phone frame not to physical scale.</div>
</div>"""


def slide_newsletter(c):
    p = c["placements"][3]
    r = c["reach"]
    subs = fmt_int(r["newsletter_subscribers"]) or chip()
    opens = f'{r["newsletter_open_rate_pct"]}%' if r["newsletter_open_rate_pct"] is not None else chip()
    ctr = f'{r["newsletter_ctr_pct"]}%' if r["newsletter_ctr_pct"] is not None else chip()
    return f"""
<div class="slide">
  {head("Placement 04", "Newsletter sponsorship", "One sponsor per issue, placed above the first editorial item. No competing advertisers in the same send.", c["brand"])}
  <div class="sbody">
    <div class="frame" style="flex:1;">
      <div class="chrome"><i></i><i></i><i></i>
        <div class="url">Inbox — The Dispatch from destination.com</div></div>
      <div class="mail">
        <div class="card">
          <div class="mh"><div class="t">The Dispatch</div>
            <div class="d">Issue 118 · Thursday 28 August</div></div>
          <div style="padding:14px 17px 4px;">
            <div class="slot" style="height:104px;margin:0;">
              <span class="tag">Sponsor placement</span>
              <span class="sz">600 × 200<small>Banner + 40 words + CTA link</small></span></div>
            <div style="font-size:8px;color:var(--ink-light);letter-spacing:1.1px;
                 text-transform:uppercase;margin-top:8px;text-align:center;">
              Presented by our sponsor</div>
          </div>
          <div class="item" style="margin-top:8px;"><div class="it">The Amalfi Coast road rule
            nobody tells you about</div>
            <div class="ix">Rental cars are barred from the SS163 on alternating days all
              summer. Here is how to plan around it.</div></div>
          <div class="item"><div class="it">Five shoulder-season trips worth booking now</div>
            <div class="ix">Late September and early October pricing, across five destinations
              our readers searched most this month.</div></div>
          <div class="item" style="border-bottom:none;"><div class="it">Points desk: the transfer
            bonus worth acting on</div>
            <div class="ix">A time-limited transfer bonus, and whether it beats booking cash.</div></div>
        </div>
      </div>
    </div>
    <div class="spec">
      <div class="card hi"><div class="ck">Unit</div><div class="cv">{p['name']}</div></div>
      <div class="card"><dl>
        <div><dt>Creative</dt><dd>{p['sizes']}</dd></div>
        <div><dt>Position</dt><dd>{p['position']}</dd></div>
        <div><dt>Exclusivity</dt><dd>{p['note']}</dd></div>
      </dl></div>
      <div class="card"><dl>
        <div><dt>List size</dt><dd>{subs}</dd></div>
        <div><dt>Open rate</dt><dd>{opens}</dd></div>
        <div><dt>Click rate</dt><dd>{ctr}</dd></div>
      </dl></div>
    </div>
  </div>
  <div class="sfoot">Mockup for placement illustration. Editorial items shown are representative
    of a typical issue.</div>
</div>"""


def slide_reach(c):
    r, g = c["reach"], c["geography"]
    return f"""
<div class="slide">
  {banner_if_missing(r["monthly_readers"], r["monthly_pageviews"], r["newsletter_subscribers"],
                     r["pages_per_session"], g)}
  {head("Audience 01", "Reach and engagement", "Newsletter and site scale, and how deeply the audience reads.", c["brand"], estimate_pill(c))}
  <div class="sbody" style="flex-direction:column;gap:26px;">
    <div class="tiles">
      {tile("Monthly readers", r["monthly_readers"])}
      {tile("Monthly pageviews", r["monthly_pageviews"])}
      {tile("Newsletter subscribers", r["newsletter_subscribers"])}
      {tile("Pages per session", r["pages_per_session"])}
    </div>
    <div style="display:flex;gap:34px;flex:1;min-height:0;">
      <div style="flex:1;display:flex;flex-direction:column;">{bars("Where readers are", g, fill=True)}</div>
      <div style="width:392px;flex-shrink:0;">
        <div class="card" style="background:var(--cream);border:1px solid var(--border);
             border-radius:12px;padding:19px 21px;height:100%;display:flex;flex-direction:column;">
          <div class="ck" style="font-size:10px;letter-spacing:1.7px;text-transform:uppercase;
               color:var(--accent);font-weight:700;margin-bottom:9px;">Engagement</div>
          <dl style="display:flex;flex-direction:column;justify-content:center;gap:24px;flex:1;">
            <div><dt style="font-size:9.5px;letter-spacing:1.5px;text-transform:uppercase;
                 color:var(--ink-light);font-weight:700;margin-bottom:3px;">Average engaged time</dt>
              <dd style="font-size:13px;">{(str(r['avg_engaged_time_min']) + ' minutes') if r['avg_engaged_time_min'] is not None else chip()}</dd></div>
            <div><dt style="font-size:9.5px;letter-spacing:1.5px;text-transform:uppercase;
                 color:var(--ink-light);font-weight:700;margin-bottom:3px;">Returning readers</dt>
              <dd style="font-size:13px;">{(str(r['returning_reader_pct']) + '%') if r['returning_reader_pct'] is not None else chip()}</dd></div>
            <div><dt style="font-size:9.5px;letter-spacing:1.5px;text-transform:uppercase;
                 color:var(--ink-light);font-weight:700;margin-bottom:3px;">Newsletter open rate</dt>
              <dd style="font-size:13px;">{(str(r['newsletter_open_rate_pct']) + '%') if r['newsletter_open_rate_pct'] is not None else chip()}</dd></div>
          </dl>
        </div>
      </div>
    </div>
  </div>
  <div class="sfoot">{c['basis_note']}</div>
</div>"""


def slide_demographics(c):
    a, gen, dev, intent = c["age"], c["gender"], c["device"], c["intent"]
    rows = []
    for it in intent:
        if it["value"] is None:
            v = chip()
        else:
            v = f'{fmt_int(it["value"])}{(" " + it["unit"]) if it["unit"] == "min" else it["unit"]}'
        rows.append(
            f'<div style="display:flex;justify-content:space-between;align-items:center;gap:16px;'
            f'padding:11px 0;border-bottom:1px solid var(--border);">'
            f'<span style="font-size:12.5px;color:var(--ink-mid);line-height:1.45;">{it["label"]}</span>'
            f'<span style="font-size:15px;font-weight:700;color:var(--primary);white-space:nowrap;">{v}</span></div>'
        )
    return f"""
<div class="slide">
  {banner_if_missing(a, gen, dev, intent)}
  {head("Audience 02", "Who is reading", "Age, gender and device against the GA4 taxonomy; intent from on-site behaviour.", c["brand"], estimate_pill(c))}
  <div class="sbody">
    <div style="flex:1;display:flex;flex-direction:column;">{bars("Age", a, fill=True)}</div>
    <div style="width:560px;flex-shrink:0;display:flex;flex-direction:column;gap:22px;">
      {split_bar("Gender", gen)}
      {split_bar("Device", dev)}
      <div style="flex:1;display:flex;flex-direction:column;min-height:0;">
        <div class="ct" style="font-size:11px;letter-spacing:1.6px;text-transform:uppercase;
             color:var(--ink-light);font-weight:700;margin-bottom:4px;">Intent signals</div>
        <div style="flex:1;display:flex;flex-direction:column;justify-content:space-evenly;">
          {"".join(rows)}
        </div>
      </div>
    </div>
  </div>
  <div class="sfoot">{c['basis_note']}</div>
</div>"""


def slide_segments(c):
    return f"""
<div class="slide">
  {head("Audience 03", "Interest and income segments",
        "How the audience maps onto the affinity and in-market segments a travel buyer plans against.",
        c["brand"], estimate_pill(c))}
  <div class="sbody">
    <div style="flex:1;display:flex;flex-direction:column;">
      {bars("Affinity &amp; in-market segments — share of audience", c["segments"],
             fill=True, label_width=232)}
    </div>
    <div style="width:520px;flex-shrink:0;display:flex;flex-direction:column;">
      {bars("Household income decile (US)", c["income"], fill=True)}
    </div>
  </div>
  <div class="sfoot">{c['segments_note']}</div>
</div>"""


def slide_specs(c):
    b = c["brand"]
    trs = "".join(
        f'<tr><td class="nm">{p["name"]}</td><td class="sz">{p["sizes"]}</td>'
        f'<td>{p["position"]}</td><td>{p["note"]}</td></tr>'
        for p in c["placements"]
    )
    return f"""
<div class="slide">
  {head("Summary", "Placement specifications", "Every unit available, in one table.", c["brand"])}
  <div class="sbody" style="flex-direction:column;">
    <table class="spectbl">
      <tr><th>Unit</th><th>Sizes / format</th><th>Position</th><th>Notes</th></tr>
      {trs}
    </table>
    <div style="margin-top:auto;display:flex;gap:16px;">
      <div class="tile" style="background:#FCF4EE;border-color:#E8C4AA;">
        <div class="k" style="color:var(--accent);">Technical</div>
        <div style="font-size:13px;color:var(--ink-mid);line-height:1.6;margin-top:10px;">
          Standard IAB sizes, third-party tags accepted. Viewability and delivery reporting
          shared at whatever cadence suits you.</div></div>
      <div class="tile">
        <div class="k">Brand safety</div>
        <div style="font-size:13px;color:var(--ink-mid);line-height:1.6;margin-top:10px;">
          Editorial travel content only. No auto-refresh, no interstitials, one in-view unit
          per viewport.</div></div>
      <div class="tile">
        <div class="k">Next step</div>
        <div style="font-size:13px;color:var(--ink-mid);line-height:1.6;margin-top:10px;">
          {b['contact_name']}<br>{b['contact_email']}</div></div>
    </div>
  </div>
  <div class="sfoot">{b['property']} · Prepared for {b['prepared_for']} · {b['date']}</div>
</div>"""


SLIDES = [
    ("01-cover", slide_cover),
    ("02-desktop-in-article", slide_in_article),
    ("03-where-to-stay-native", slide_native),
    ("04-mobile-in-feed", slide_mobile),
    ("05-newsletter-sponsorship", slide_newsletter),
    ("06-audience-reach", slide_reach),
    ("07-audience-demographics", slide_demographics),
    ("08-audience-segments", slide_segments),
    ("09-placement-specs", slide_specs),
]


# ---------------------------------------------------------------- render

def render(name, body):
    html_path = os.path.join(WORK, name + ".html")
    png_path = os.path.join(WORK, name + ".png")
    jpg_path = os.path.join(OUT, name + ".jpg")

    with open(FONT_CSS) as f:
        fonts = f.read()
    doc = ("<!doctype html><html><head><meta charset='utf-8'>"
           f"<style>{fonts}</style><style>{CSS}</style></head><body>{body}</body></html>")
    with open(html_path, "w") as f:
        f.write(doc)

    if os.path.exists(png_path):
        os.remove(png_path)
    args = [CHROME]
    if not IS_SHELL:
        args.append("--headless=new")
    args += ["--no-sandbox", "--disable-gpu", "--hide-scrollbars",
             "--disable-dev-shm-usage", "--force-color-profile=srgb",
             f"--force-device-scale-factor={SCALE}",
             f"--window-size={WIDTH},{HEIGHT}", "--virtual-time-budget=4000",
             f"--screenshot={png_path}", "file://" + html_path]
    subprocess.run(args, check=True, capture_output=True)
    if not os.path.exists(png_path):
        raise RuntimeError(f"Chromium produced no screenshot for {name}")

    im = Image.open(png_path).convert("RGB")
    if im.size != (WIDTH * SCALE, HEIGHT * SCALE):
        raise RuntimeError(
            f"{name}: expected {WIDTH*SCALE}x{HEIGHT*SCALE} capture, got {im.width}x{im.height}. "
            "The browser viewport does not match the slide size."
        )
    if im.width != FINAL_WIDTH:
        im = im.resize((FINAL_WIDTH, round(im.height * FINAL_WIDTH / im.width)), Image.LANCZOS)
    im.save(jpg_path, "JPEG", quality=JPG_QUALITY, optimize=True, progressive=True)
    return jpg_path, im.size, os.path.getsize(jpg_path)


# ---------------------------------------------------------------- deck export

DECK_BASENAME = "destination-com-media-kit"


def build_deck(jpgs, cfg):
    """Assemble the rendered slides into a PDF and a 16:9 PPTX."""
    made = []

    # PDF: 2400x1350 at 150dpi is exactly 16x9in, so pages come out full-bleed
    # with no scaling and no margin.
    pdf_path = os.path.join(OUT, DECK_BASENAME + ".pdf")
    pages = [Image.open(j).convert("RGB") for j in jpgs]
    pages[0].save(pdf_path, "PDF", save_all=True, append_images=pages[1:],
                  resolution=150.0, title=f"{cfg['brand']['property']} media kit")
    for im in pages:
        im.close()
    made.append(pdf_path)

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("  (skipped .pptx — pip install python-pptx)")
        return made

    pptx_path = os.path.join(OUT, DECK_BASENAME + ".pptx")
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    for j in jpgs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(j, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.core_properties.title = f"{cfg['brand']['property']} media kit"
    prs.core_properties.author = cfg["brand"]["contact_name"]
    # PPTX core properties cap each field at 255 chars, so the full basis_note
    # does not fit here — it is on the slides themselves, which is where it counts.
    comment = f"Prepared for {cfg['brand']['prepared_for']}, {cfg['brand']['date']}."
    if cfg.get("figures_basis") == "estimated":
        comment += (" Audience figures are internal estimates pending GA4/ESP"
                    " verification — see the note on each audience slide.")
    prs.core_properties.comments = comment[:255]
    prs.save(pptx_path)
    made.append(pptx_path)
    return made


def main():
    os.makedirs(WORK, exist_ok=True)
    os.makedirs(OUT, exist_ok=True)
    cfg = load_cfg()
    print(f"Chromium: {CHROME}\n")
    jpgs = []
    for name, fn in SLIDES:
        path, size, nbytes = render(name, fn(cfg))
        jpgs.append(path)
        print(f"  {os.path.basename(path):38s} {size[0]}×{size[1]}  {nbytes/1024:6.0f} KB")
    print(f"\nWrote {len(jpgs)} JPGs to {OUT}")

    for d in build_deck(jpgs, cfg):
        print(f"  {os.path.basename(d):38s} {os.path.getsize(d)/1024:6.0f} KB")

    missing = [k for k, v in cfg["reach"].items() if v is None]
    if missing:
        print(f"\nWARNING: {len(missing)} reach figures still null "
              f"({', '.join(missing)}) — those slides carry the TEMPLATE banner.")
    if cfg.get("figures_basis") == "estimated":
        print("\nNOTE: figures_basis is \"estimated\" — every data slide carries the "
              "ESTIMATES pill and prints basis_note in its footer.\n"
              "      Set it to \"measured\" only once the figures are real.")


if __name__ == "__main__":
    main()
